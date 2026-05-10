"""
Local RAG chatbot — no frameworks, no cloud.

Usage:
    python rag.py doc1.pdf doc2.txt          # load docs and start chat
    python rag.py                            # start chat using saved index

Dependencies:
    pip install faiss-cpu sentence-transformers pymupdf requests
    pip install python-docx python-pptx pytesseract pillow   # optional formats
    brew install tesseract                                    # scanned PDFs only
"""
from __future__ import annotations

import sys
import time
import pickle
import textwrap
import argparse
from pathlib import Path
from typing import Optional
from contextlib import contextmanager
from dataclasses import dataclass, field

import numpy as np
import faiss
import requests
from sentence_transformers import SentenceTransformer

# ─── Config ──────────────────────────────────────────────────────────────────

INDEX_PATH      = "rag_index"
EMBED_MODEL     = "all-MiniLM-L6-v2"
CHUNK_SIZE      = 400
CHUNK_OVERLAP   = 50
TOP_K           = 5
HISTORY_TURNS   = 4       # past Q&A pairs included in each prompt
SCORE_THRESHOLD = 0.25    # min cosine similarity to accept a retrieval result

OLLAMA_URL   = "http://localhost:11434/api/generate"
OLLAMA_MODEL = "mistral"

# ─── Timing ──────────────────────────────────────────────────────────────────

@contextmanager
def timer(label: str):
    t0 = time.perf_counter()
    yield
    print(f"  [{(time.perf_counter() - t0) * 1000:.0f} ms] {label}")

# ─── Document Loaders ────────────────────────────────────────────────────────

def _load_pdf(p: Path) -> str:
    try:
        import fitz
    except ImportError:
        sys.exit("pip install pymupdf")
    doc = fitz.open(str(p))
    pages = []
    for page in doc:
        text = page.get_text().strip()
        pages.append(text if text else _ocr_page(page))
    return "\n".join(pages)

def _ocr_page(page) -> str:
    try:
        import pytesseract
        from PIL import Image
        import io
    except ImportError:
        sys.exit("Scanned PDF needs: pip install pytesseract pillow && brew install tesseract")
    pix = page.get_pixmap(dpi=300)
    return pytesseract.image_to_string(Image.open(io.BytesIO(pix.tobytes("png"))))

def _load_docx(p: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        sys.exit("pip install python-docx")
    return "\n".join(para.text for para in Document(str(p)).paragraphs if para.text.strip())

def _load_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("pip install python-pptx")
    slides = []
    for slide in Presentation(str(p)).slides:
        text = " ".join(s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip())
        if text.strip():
            slides.append(text)
    return "\n".join(slides)

def load_file(path: str) -> str:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(path)
    ext = p.suffix.lower()
    if ext == ".pdf":  return _load_pdf(p)
    if ext == ".docx": return _load_docx(p)
    if ext == ".pptx": return _load_pptx(p)
    return p.read_text(encoding="utf-8", errors="replace")

# ─── Chunking ────────────────────────────────────────────────────────────────

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    step  = max(1, size - overlap)
    return [
        " ".join(words[i : i + size])
        for i in range(0, len(words), step)
        if words[i : i + size]
    ]

# ─── Embedder (singleton) ────────────────────────────────────────────────────

_embedder: Optional[SentenceTransformer] = None

def get_embedder() -> SentenceTransformer:
    global _embedder
    if _embedder is None:
        print(f"Loading embedding model '{EMBED_MODEL}'...")
        _embedder = SentenceTransformer(EMBED_MODEL)
    return _embedder

def embed(texts: list[str]) -> np.ndarray:
    vecs = get_embedder().encode(texts, normalize_embeddings=True, show_progress_bar=False)
    return np.array(vecs, dtype="float32")

# ─── In-memory Index ─────────────────────────────────────────────────────────

@dataclass
class Index:
    faiss_index: faiss.Index
    chunks: list[str]

    @classmethod
    def build(cls, chunks: list[str], vecs: np.ndarray) -> Index:
        idx = faiss.IndexFlatIP(vecs.shape[1])
        idx.add(vecs)
        return cls(idx, chunks)

    @classmethod
    def load_from_disk(cls) -> Index:
        fp = f"{INDEX_PATH}.faiss"
        if not Path(fp).exists():
            sys.exit(f"No saved index at '{fp}'. Pass document files to load them.")
        idx    = faiss.read_index(fp)
        chunks = pickle.loads(Path(f"{INDEX_PATH}.pkl").read_bytes())
        return cls(idx, chunks)

    def save(self) -> None:
        faiss.write_index(self.faiss_index, f"{INDEX_PATH}.faiss")
        Path(f"{INDEX_PATH}.pkl").write_bytes(pickle.dumps(self.chunks))

    def search(self, query: str, k: int = TOP_K) -> list[tuple[float, str]]:
        q_vec = embed([query])
        scores, indices = self.faiss_index.search(q_vec, k)
        return [
            (float(s), self.chunks[i])
            for s, i in zip(scores[0], indices[0])
            if i != -1
        ]

# ─── LLM ─────────────────────────────────────────────────────────────────────

def call_llm(prompt: str) -> str:
    try:
        resp = requests.post(
            OLLAMA_URL,
            json={"model": OLLAMA_MODEL, "prompt": prompt, "stream": False},
            timeout=180,
        )
        resp.raise_for_status()
        return resp.json().get("response", "").strip()
    except requests.exceptions.ConnectionError:
        sys.exit(f"\nCannot reach Ollama at {OLLAMA_URL}.\nRun: ollama serve && ollama pull {OLLAMA_MODEL}\n")

# ─── Prompt ──────────────────────────────────────────────────────────────────

@dataclass
class Turn:
    question: str
    answer: str

def build_prompt(question: str, context_chunks: list[str], history: list[Turn]) -> str:
    context = "\n\n---\n\n".join(context_chunks)

    history_block = ""
    if history:
        lines = []
        for t in history[-HISTORY_TURNS:]:
            lines.append(f"User: {t.question}")
            lines.append(f"Assistant: {t.answer}")
        history_block = "\nCONVERSATION HISTORY:\n" + "\n".join(lines) + "\n"

    return textwrap.dedent(f"""\
        You are a helpful assistant. Answer using only the context provided.
        If the answer is not in the context, say "I don't know based on the provided documents."
        Use the conversation history to resolve follow-up questions and references to prior answers.

        CONTEXT:
        {context}
        {history_block}
        CURRENT QUESTION:
        {question}

        ANSWER:
    """)

# ─── Chat Loop ───────────────────────────────────────────────────────────────

HELP = """
  Commands:
    history   — show all Q&A from this session
    clear     — clear conversation history
    quit      — exit
"""

def chat(index: Index) -> None:
    history: list[Turn] = []
    print("\nReady. Ask a question about your documents.")
    print(HELP)

    while True:
        try:
            question = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye.")
            break

        if not question:
            continue

        if question.lower() in {"quit", "exit", "q"}:
            print("Bye.")
            break

        if question.lower() == "history":
            if not history:
                print("  No questions asked yet.\n")
            else:
                for i, t in enumerate(history, 1):
                    print(f"  Q{i}: {t.question}")
                    preview = t.answer[:160] + ("..." if len(t.answer) > 160 else "")
                    print(f"  A{i}: {preview}\n")
            continue

        if question.lower() == "clear":
            history.clear()
            print("  Conversation history cleared.\n")
            continue

        with timer("retrieve"):
            hits = index.search(question, k=TOP_K)

        if not hits or hits[0][0] < SCORE_THRESHOLD:
            print("  Bot: I couldn't find relevant information in the documents for that question.\n")
            continue

        context_chunks = [chunk for _, chunk in hits]
        prompt = build_prompt(question, context_chunks, history)

        with timer("llm"):
            answer = call_llm(prompt)

        print(f"\n  Bot: {answer}\n")
        history.append(Turn(question=question, answer=answer))

# ─── Entry Point ─────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Local RAG chatbot — ask questions about your documents",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
            Examples:
              python rag.py report.pdf notes.txt    # load docs and chat
              python rag.py                         # chat using saved index
        """),
    )
    parser.add_argument("files", nargs="*", help="PDF, DOCX, PPTX, or TXT files to load")
    args = parser.parse_args()

    if args.files:
        all_chunks: list[str] = []
        for path in args.files:
            print(f"Loading  {path}")
            with timer("load + chunk"):
                text   = load_file(path)
                chunks = chunk_text(text)
            print(f"  → {len(chunks)} chunks")
            all_chunks.extend(chunks)

        if not all_chunks:
            sys.exit("No content extracted. Check your files.")

        print(f"\nEmbedding {len(all_chunks)} chunks...")
        with timer("embed"):
            vecs = embed(all_chunks)

        with timer("build index"):
            index = Index.build(all_chunks, vecs)

        index.save()
        print(f"Index ready — {index.faiss_index.ntotal} vectors.\n")
    else:
        print("No files given — loading saved index from disk...")
        with timer("load index"):
            index = Index.load_from_disk()
        print(f"Loaded {index.faiss_index.ntotal} vectors.\n")

    chat(index)

if __name__ == "__main__":
    main()
